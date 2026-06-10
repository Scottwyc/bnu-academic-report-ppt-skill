#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


SKILL_DIR = Path(__file__).resolve().parents[1]
DEFAULT_ASSETS = SKILL_DIR / "assets"


def default_academic_builder() -> Path:
    env_path = os.environ.get("ACADEMIC_PPT_BUILDER_SCRIPT")
    if env_path:
        return Path(env_path).expanduser()

    codex_home = Path(os.environ.get("CODEX_HOME", "~/.codex")).expanduser()
    candidates = [
        SKILL_DIR.parent / "academic-ppt-builder" / "scripts" / "render_academic_pptx.py",
        codex_home / "skills" / "general" / "academic-ppt-builder" / "scripts" / "render_academic_pptx.py",
        Path("~/.codex/skills/general/academic-ppt-builder/scripts/render_academic_pptx.py").expanduser(),
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0]


def run(cmd: list[str], cwd: Path | None = None) -> None:
    print("+", " ".join(cmd), flush=True)
    subprocess.run(cmd, cwd=str(cwd) if cwd else None, check=True)


def load_outline(path: Path) -> list[dict]:
    raw = json.loads(path.read_text(encoding="utf-8"))
    slides = raw.get("slides", raw) if isinstance(raw, dict) else raw
    return [s for s in slides if isinstance(s, dict)]


def outline_image_paths(slide: dict) -> list[str]:
    paths: list[str] = []
    for key in ("image_path", "image_paths", "images"):
        value = slide.get(key)
        if isinstance(value, str):
            paths.append(value)
        elif isinstance(value, list):
            for item in value:
                if isinstance(item, str):
                    paths.append(item)
                elif isinstance(item, dict) and item.get("path"):
                    paths.append(str(item["path"]))
    return paths


def path_exists(path_text: str, md_dir: Path, ppt_dir: Path) -> bool:
    p = Path(path_text)
    candidates = []
    if p.is_absolute():
        candidates.append(p)
    else:
        candidates.extend([md_dir / p, ppt_dir / p, Path.cwd() / p])
    return any(c.exists() for c in candidates)


def validate_pptx(pptx: Path, outline_json: Path, input_md: Path) -> dict:
    result: dict = {}
    prs = Presentation(pptx)
    result["slides"] = len(prs.slides)
    result["pptx_size"] = pptx.stat().st_size
    result["outline_size"] = outline_json.stat().st_size if outline_json.exists() else 0

    missing: list[tuple[int, str]] = []
    if outline_json.exists():
        for idx, slide in enumerate(load_outline(outline_json), 1):
            for img in outline_image_paths(slide):
                if img and not path_exists(img, input_md.parent, pptx.parent):
                    missing.append((idx, img))
    result["missing_images"] = missing

    max_dev = 0.0
    aspect_issues: list[tuple[int, float, str]] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type != 13:
                continue
            img = shape.image
            orig = img.size[0] / img.size[1]
            disp = shape.width / shape.height
            dev = abs(disp / orig - 1)
            max_dev = max(max_dev, dev)
            if dev > 0.03:
                aspect_issues.append((slide_idx, round(dev, 4), img.filename))
    result["aspect_max_dev"] = round(max_dev, 5)
    result["aspect_issues"] = aspect_issues

    with zipfile.ZipFile(pptx) as zf:
        result["pptx_zip_ok"] = zf.testzip() is None
    return result


def export_preview(pptx: Path, preview_dir: Path, dpi: int) -> Path | None:
    if not shutil.which("libreoffice"):
        print("Preview skipped: libreoffice not found.", file=sys.stderr)
        return None
    preview_dir.mkdir(parents=True, exist_ok=True)
    run(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", str(preview_dir), str(pptx)])
    pdf = preview_dir / f"{pptx.stem}.pdf"
    if shutil.which("pdftoppm") and pdf.exists():
        run(["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(preview_dir / "slide")])
    return pdf if pdf.exists() else None


def main() -> None:
    parser = argparse.ArgumentParser(description="Render and validate a BNU academic report style PPTX.")
    parser.add_argument("--input-md", required=True, type=Path, help="PPT-ready Markdown file.")
    parser.add_argument("--output-pptx", required=True, type=Path, help="Output editable PPTX.")
    parser.add_argument("--outline-json", type=Path, help="Output slide outline JSON.")
    parser.add_argument("--preview-dir", type=Path, help="Directory for PDF/PNG preview export.")
    parser.add_argument("--preview-dpi", type=int, default=160)
    parser.add_argument("--builder-script", type=Path, default=default_academic_builder())
    parser.add_argument("--assets-dir", type=Path, default=DEFAULT_ASSETS)
    parser.add_argument("--no-preview", action="store_true")
    args = parser.parse_args()

    input_md = args.input_md.resolve()
    output_pptx = args.output_pptx.resolve()
    outline_json = (args.outline_json or output_pptx.with_name(output_pptx.stem + "__slide_outline.json")).resolve()
    builder = args.builder_script.resolve()
    assets_dir = args.assets_dir.resolve()

    if not input_md.exists():
        raise FileNotFoundError(input_md)
    if not builder.exists():
        raise FileNotFoundError(builder)
    if not assets_dir.exists():
        raise FileNotFoundError(assets_dir)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    outline_json.parent.mkdir(parents=True, exist_ok=True)

    run([
        sys.executable,
        str(builder),
        str(input_md),
        str(output_pptx),
        "--outline-json",
        str(outline_json),
        "--assets-dir",
        str(assets_dir),
    ])
    validation = validate_pptx(output_pptx, outline_json, input_md)

    pdf = None
    if not args.no_preview:
        preview_dir = (args.preview_dir or output_pptx.parent / f"preview_{output_pptx.stem}").resolve()
        pdf = export_preview(output_pptx, preview_dir, args.preview_dpi)
        validation["preview_pdf"] = str(pdf) if pdf else None

    print(json.dumps({
        "output_pptx": str(output_pptx),
        "outline_json": str(outline_json),
        "validation": validation,
    }, ensure_ascii=False, indent=2))

    if validation["missing_images"] or validation["aspect_issues"] or not validation["pptx_zip_ok"]:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
